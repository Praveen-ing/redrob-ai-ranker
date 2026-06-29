"""
Generate a .pptx PowerPoint presentation for Redrob AI Ranker — INDIA.RUNS Hackathon
Matches the official template: dark header bar, purple INDIA.RUNS branding, gradient accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────── COLORS ────────────────
BLACK       = RGBColor(0x0A, 0x0A, 0x0A)
DARK        = RGBColor(0x11, 0x11, 0x11)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_400    = RGBColor(0xA3, 0xA3, 0xA3)
GRAY_500    = RGBColor(0x73, 0x73, 0x73)
GRAY_700    = RGBColor(0x40, 0x40, 0x40)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_DARK = RGBColor(0x6D, 0x28, 0xD9)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)
GREEN       = RGBColor(0x05, 0x96, 0x69)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
RED         = RGBColor(0xEF, 0x44, 0x44)
LIGHT_BG    = RGBColor(0xF9, 0xF8, 0xFF)
ACCENT_1    = RGBColor(0x63, 0x66, 0xF1)  # indigo
ACCENT_2    = RGBColor(0xA8, 0x55, 0xF7)  # violet

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use the blank layout
blank_layout = prs.slide_layouts[6]

# ──────────────── HELPERS ────────────────

def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(4), font_name='Calibri'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_rich_paragraph(text_frame, runs_list, font_size=14, alignment=PP_ALIGN.LEFT,
                       space_before=Pt(0), space_after=Pt(4), font_name='Calibri'):
    """Add a paragraph with multiple runs (bold/normal mixed)."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    for (text, bold, color) in runs_list:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return p

def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=13,
                       color=RGBColor(0x1E,0x1E,0x1E), bold_prefix=True, font_name='Calibri'):
    """Add a textbox with bullet points. Each bullet can be (bold_part, rest) or just a string."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        # Add bullet character
        if isinstance(bullet, tuple):
            bold_text, normal_text = bullet
            run_b = p.add_run()
            run_b.text = "●  " + bold_text
            run_b.font.size = Pt(font_size)
            run_b.font.bold = True
            run_b.font.color.rgb = BLACK
            run_b.font.name = font_name
            if normal_text:
                run_n = p.add_run()
                run_n.text = " " + normal_text
                run_n.font.size = Pt(font_size)
                run_n.font.bold = False
                run_n.font.color.rgb = color
                run_n.font.name = font_name
        else:
            run = p.add_run()
            run.text = "●  " + bullet
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = font_name
    return txBox

def add_dark_header(slide):
    """Add the dark header bar with 'redrob | H2S' and 'INDIA.RUNS'."""
    # Dark background bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.65), DARK)

    # redrob logo text
    tb = add_textbox(slide, Inches(0.5), Inches(0.12), Inches(2), Inches(0.4),
                     "● redrob  │  H2S", font_size=12, bold=True, color=WHITE)

    # INDIA.RUNS on the right
    add_textbox(slide, Inches(10.5), Inches(0.1), Inches(2.5), Inches(0.45),
                "INDIA.RUNS", font_size=15, bold=True, color=PURPLE,
                alignment=PP_ALIGN.RIGHT)

def add_accent_bar(slide):
    """Add the gradient-like accent bar at the bottom."""
    bar_h = Inches(0.06)
    bar_y = SLIDE_H - bar_h
    # We simulate gradient with 3 thin color strips
    third = Inches(13.333 / 3)
    add_shape(slide, Inches(0), bar_y, third, bar_h, ACCENT_1)
    add_shape(slide, third, bar_y, third, bar_h, PURPLE)
    add_shape(slide, third * 2, bar_y, third, bar_h, ACCENT_2)

def add_slide_title(slide, title, subtitle=None):
    """Add slide title text below the header."""
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.55),
                title, font_size=24, bold=True, color=BLACK)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.35),
                    subtitle, font_size=10, bold=True, color=GRAY_500)

def make_content_slide(title, subtitle=None):
    """Create a standard content slide with header, title, and accent bar."""
    slide = prs.slides.add_slide(blank_layout)
    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_dark_header(slide)
    add_accent_bar(slide)
    add_slide_title(slide, title, subtitle)
    return slide


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# Hero gradient background (top half) — we use a solid dark color since pptx doesn't do CSS gradients
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(4.2), RGBColor(0x14, 0x0A, 0x2E))
# Add subtle overlay rectangles to simulate gradient
add_shape(slide1, Inches(0), Inches(0), Inches(5), Inches(4.2), RGBColor(0x1A, 0x05, 0x33))
add_shape(slide1, Inches(8), Inches(1), Inches(5.333), Inches(3.2), RGBColor(0x4A, 0x19, 0x42))

# Logos
add_textbox(slide1, Inches(4), Inches(0.5), Inches(5.333), Inches(0.5),
            "● redrob  │  H2S", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# INDIA.RUNS title
add_textbox(slide1, Inches(2), Inches(1.1), Inches(9.333), Inches(1.2),
            "INDIA.RUNS", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tagline pill
tagline_box = add_textbox(slide1, Inches(4), Inches(2.6), Inches(5.333), Inches(0.5),
                          "Build what next India runs on", font_size=13, bold=False,
                          color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# White bottom section
add_shape(slide1, Inches(0), Inches(4.2), SLIDE_W, Inches(3.3), WHITE)

# Team info
info_items = [
    ("Team Name : ", "TOPGUN"),
    ("Team Leader Name : ", "Nethavath Praveen"),
    ("Problem Statement : ", "Build an intelligent candidate discovery & ranking system that goes beyond keyword matching to evaluate 100,000+ profiles using semantic understanding, career trajectory analysis, and behavioral intent scoring — fully offline, in under 5 minutes."),
]
y = 4.5
for label, value in info_items:
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = label
    run_b.font.size = Pt(14)
    run_b.font.bold = True
    run_b.font.color.rgb = BLACK
    run_b.font.name = 'Calibri'
    run_v = p.add_run()
    run_v.text = value
    run_v.font.size = Pt(14)
    run_v.font.bold = False
    run_v.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    run_v.font.name = 'Calibri'
    y += 0.55 if label != "Problem Statement : " else 0.9

# Accent bar
add_accent_bar(slide1)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════
slide2 = make_content_slide("Solution Overview")
add_bullet_textbox(slide2, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    ("Redrob AI Ranker —", "A production-grade, fully offline AI candidate ranking engine that processes 100,000+ profiles in ~34 seconds with zero API calls, zero internet dependency, and fully transparent scoring."),
    ("Semantic Understanding:", "Instead of counting keywords, skills are grouped into 5 semantic clusters (Core Programming, Vector DBs, Retrieval & Ranking, LLM Engineering, Infrastructure). Candidates get credit for the entire cluster they demonstrate competency in."),
    ("What differentiates us:", "We use TF-IDF weighted cluster matching, Needleman-Wunsch career alignment (from bioinformatics), PageRank-based corporate pedigree scoring, Shannon entropy anomaly detection, and a recruiter active learning loop — all running fully offline with Python's standard library."),
    ("Every ranking decision is explainable —", "deterministic, data-driven XAI rationales are generated for each shortlisted candidate, not generic templates."),
], font_size=13)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — JD UNDERSTANDING & CANDIDATE EVALUATION
# ══════════════════════════════════════════════════════════════════
slide3 = make_content_slide("JD Understanding & Candidate Evaluation")
add_bullet_textbox(slide3, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    ("Key requirements extracted from the JD", "are codified in a single JSON config (jd_rules.json): target role (Senior AI Engineer), ideal experience band (5–9 yrs), 5 semantic skill clusters with 50+ accepted technologies, preferred/acceptable locations, and company tier classifications."),
    ("Beyond keyword matching:", "We evaluate candidate fit using IDF-weighted semantic cluster coverage, skill endorsement bonuses (up to 15%), platform assessment score integration, profile summary & headline text mining, and a cohesion check that discounts advanced skills without foundational programming."),
    ("Full signal utilization:", "Every field in the candidate schema is used — education tier, skill assessment scores, recruiter response rate, notice period, interview completion rate, offer acceptance rate, saved-by-recruiters count, and response time — not just skills and experience."),
    ("Anti-gaming integrity:", "Honeypot detection filters out fake profiles using impossible skill claims, assessment contradictions, timeline inconsistencies, and Shannon entropy anomaly detection."),
], font_size=13)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — RANKING METHODOLOGY
# ══════════════════════════════════════════════════════════════════
slide4 = make_content_slide("Ranking Methodology", "4-AXIS WEIGHTED COMPOSITE SCORING")

pillar_data = [
    ("🧠 Semantic Competency — 40%",
     "TF-IDF weighted cluster matching across 5 skill domains. Endorsement bonuses, assessment score integration, and cohesion penalty for skills without foundational programming.",
     RGBColor(0xF5, 0xF3, 0xFF), PURPLE_DARK),
    ("📈 Career Trajectory — 30%",
     "Needleman-Wunsch career sequence alignment, PageRank corporate pedigree, promotion velocity gradient, education tier scoring, and consulting penalty.",
     RGBColor(0xEF, 0xF6, 0xFF), BLUE),
    ("🎯 Behavioral Intent — 20%",
     "Recruiter response rate, open-to-work flag, notice period, interview completion rate, offer acceptance rate, saved-by-recruiters, response time, and exponential activity decay.",
     RGBColor(0xEC, 0xFD, 0xF5), GREEN),
    ("📍 Location Intelligence — 10%",
     "Preferred cities (full score), acceptable cities (0.7), willing-to-relocate (0.6), and remote/other (0.2 baseline).",
     RGBColor(0xFF, 0xF7, 0xED), AMBER),
]

positions = [(0.8, 1.7), (6.6, 1.7), (0.8, 3.5), (6.6, 3.5)]
for idx, ((title, desc, bg_color, title_color), (x, y)) in enumerate(zip(pillar_data, positions)):
    # Card background
    card = add_shape(slide4, Inches(x), Inches(y), Inches(5.5), Inches(1.6), bg_color)
    card.shadow.inherit = False

    # Card title
    add_textbox(slide4, Inches(x + 0.2), Inches(y + 0.15), Inches(5.1), Inches(0.4),
                title, font_size=12, bold=True, color=title_color)
    # Card description
    add_textbox(slide4, Inches(x + 0.2), Inches(y + 0.55), Inches(5.1), Inches(0.95),
                desc, font_size=10, bold=False, color=GRAY_700)

# Formula bar
formula_bg = add_shape(slide4, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.55), RGBColor(0x1E, 0x1E, 0x2E))
add_textbox(slide4, Inches(0.8), Inches(5.42), Inches(11.5), Inches(0.5),
            "Final Score = (Skills × 0.40) + (Experience × 0.30) + (Behavior × 0.20) + (Location × 0.10)",
            font_size=12, bold=False, color=RGBColor(0xCD, 0xD6, 0xF4), alignment=PP_ALIGN.CENTER,
            font_name='Consolas')


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — EXPLAINABILITY & DATA VALIDATION
# ══════════════════════════════════════════════════════════════════
slide5 = make_content_slide("Explainability & Data Validation")
add_bullet_textbox(slide5, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    ("How ranking decisions are explained:", 'Every shortlisted candidate receives a deterministic, context-aware XAI rationale built from actual matched clusters, career velocity metrics, and behavioral signals — not templates. Example: "Exceptional technical fit, covering 5/5 semantic clusters. 6.4 years of steady progression. Tier-1 academic background. Open to work, saved by 26 recruiters."'),
    ("Preventing hallucinations:", "The system is fully deterministic — no LLM inference, no probabilistic outputs. Every score is computed from explicit rules, every rationale is assembled from actual data points. Run it 100 times — identical output."),
    ("Handling suspicious profiles — Multi-layer honeypot detection:", '"Expert" proficiency + 0 months usage = instant disqualification. Assessment contradiction: claims expert but scored <30 on platform test. Timeline fraud: 10+ years stated, career history shows <1 year. Shannon Entropy Anomaly: text entropy deviating >2σ from population mean flags keyword-stuffed/machine-generated profiles.'),
], font_size=12)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — END-TO-END WORKFLOW
# ══════════════════════════════════════════════════════════════════
slide6 = make_content_slide("End-to-End Workflow")

# CLI Pipeline label
add_textbox(slide6, Inches(0.8), Inches(1.5), Inches(3), Inches(0.3),
            "CLI RANKING PIPELINE", font_size=9, bold=True, color=PURPLE)

# Flow boxes — Row 1 (CLI Pipeline)
flow_items_1 = [
    ("INPUT", "candidates.jsonl\n100K+ profiles", True),
    ("PHASE 1", "Pre-Scan\nTF-IDF · Entropy · PageRank", False),
    ("PHASE 2", "Parallel Scoring\nN CPU cores · 4-axis eval", False),
    ("PHASE 3", "Sort & Rank\nDeterministic tie-break", False),
    ("PHASE 4", "XAI Reasoning\nTop-100 rationales", False),
    ("OUTPUT", "team_submission.csv", True),
]
box_w = Inches(1.7)
box_h = Inches(0.9)
start_x = 0.8
gap = 0.15
for i, (label, text, is_dark) in enumerate(flow_items_1):
    x = Inches(start_x + i * 2.0)
    y = Inches(1.85)
    bg = RGBColor(0x1E, 0x1B, 0x4B) if is_dark else LIGHT_BG
    box = add_shape(slide6, x, y, box_w, box_h, bg)
    box.shadow.inherit = False

    label_color = RGBColor(0xAA, 0xAA, 0xCC) if is_dark else PURPLE
    text_color = WHITE if is_dark else BLACK

    add_textbox(slide6, x + Inches(0.05), y + Inches(0.05), box_w - Inches(0.1), Inches(0.2),
                label, font_size=7, bold=True, color=label_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide6, x + Inches(0.05), y + Inches(0.25), box_w - Inches(0.1), Inches(0.6),
                text, font_size=9, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(flow_items_1) - 1:
        add_textbox(slide6, x + box_w, y + Inches(0.25), Inches(0.3), Inches(0.35),
                    "→", font_size=14, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Dashboard flow label
add_textbox(slide6, Inches(0.8), Inches(3.1), Inches(4), Inches(0.3),
            "INTERACTIVE DASHBOARD FLOW", font_size=9, bold=True, color=PURPLE)

# Flow boxes — Row 2 (Dashboard)
flow_items_2 = [
    ("STEP 1", "Define Role"),
    ("STEP 2", "Upload CSV"),
    ("STEP 3", "Processing"),
    ("STEP 4", "Results Dashboard"),
    ("STEP 5", "Invite / Decline"),
    ("FEEDBACK", "Active Learning"),
]
for i, (label, text) in enumerate(flow_items_2):
    x = Inches(start_x + i * 2.0)
    y = Inches(3.45)
    box = add_shape(slide6, x, y, box_w, Inches(0.7), LIGHT_BG)
    box.shadow.inherit = False

    add_textbox(slide6, x + Inches(0.05), y + Inches(0.05), box_w - Inches(0.1), Inches(0.2),
                label, font_size=7, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide6, x + Inches(0.05), y + Inches(0.22), box_w - Inches(0.1), Inches(0.4),
                text, font_size=10, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

    if i < len(flow_items_2) - 1:
        add_textbox(slide6, x + box_w, y + Inches(0.12), Inches(0.3), Inches(0.35),
                    "→", font_size=14, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Key note
key_bg = add_shape(slide6, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.85), LIGHT_BG)
txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.1), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run_b = p.add_run()
run_b.text = "Key: "
run_b.font.size = Pt(10)
run_b.font.bold = True
run_b.font.color.rgb = BLACK
run_b.font.name = 'Calibri'
run_n = p.add_run()
run_n.text = "The CLI engine (top row) generates the ranked CSV. The dashboard (bottom row) consumes it and lets recruiters interactively explore, compare, and provide feedback that updates scoring weights via a perceptron learning step. All processing is local — zero network calls."
run_n.font.size = Pt(10)
run_n.font.color.rgb = GRAY_700
run_n.font.name = 'Calibri'


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════
slide7 = make_content_slide("System Architecture")

arch_data = [
    (0.8, 1.6, 5.5, "LAYER 1 — RANKING ENGINE (PYTHON CLI)", [
        "run.py — Main orchestrator. Reads JSONL/GZ, runs pipeline, outputs CSV.",
        "scorer.py — Core 4-axis scoring engine: semantic clusters, trajectory, behavior, location.",
        "reasoning_engine.py — XAI rationale generator for top-100 candidates.",
        "config.py — Loads jd_rules.json, defines scoring weights.",
    ]),
    (6.8, 1.6, 5.5, "PRE-SCAN PHASE (GLOBAL CONTEXT)", [
        "TF-IDF Calculator — Computes rarity weights across all keywords in pool.",
        "PageRank Graph — Builds company transition graph, runs 10 iterations.",
        "Shannon Entropy Baseline — Computes mean & σ for anomaly detection.",
    ]),
    (0.8, 4.0, 5.5, "LAYER 2 — INTERACTIVE DASHBOARD (HTML/JS)", [
        "Pure HTML + JavaScript single-file application.",
        "No Node.js, no build steps, no npm, no framework.",
        "Weight sliders, A/B comparison, active learning loop.",
        "LocalStorage-persisted recruiter preference profiles.",
    ]),
    (6.8, 4.0, 5.5, "ALGORITHMIC HIGHLIGHTS", [
        "Needleman-Wunsch — DP career sequence alignment (bioinformatics).",
        "PageRank — Dynamic corporate pedigree scoring.",
        "Shannon Entropy — Statistical anti-gaming detection.",
        "RLRF — Single-layer perceptron for recruiter active learning.",
    ]),
]

for (x, y, w, title, items) in arch_data:
    card_h = 2.1 if len(items) >= 4 else 1.8
    card = add_shape(slide7, Inches(x), Inches(y), Inches(w), Inches(card_h), LIGHT_BG)

    add_textbox(slide7, Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.3),
                title, font_size=9, bold=True, color=PURPLE)

    for j, item in enumerate(items):
        parts = item.split(" — ", 1) if " — " in item else [item]
        txBox = slide7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45 + j * 0.35), Inches(w - 0.4), Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if len(parts) == 2:
            run_b = p.add_run()
            run_b.text = "• " + parts[0] + " — "
            run_b.font.size = Pt(9)
            run_b.font.bold = True
            run_b.font.color.rgb = BLACK
            run_b.font.name = 'Calibri'
            run_n = p.add_run()
            run_n.text = parts[1]
            run_n.font.size = Pt(9)
            run_n.font.color.rgb = GRAY_700
            run_n.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(9)
            run.font.color.rgb = GRAY_700
            run.font.name = 'Calibri'


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS & PERFORMANCE
# ══════════════════════════════════════════════════════════════════
slide8 = make_content_slide("Results & Performance")

# Metric cards
metrics = [
    ("100,000+", "PROFILES PROCESSED"),
    ("~34s", "PROCESSING TIME"),
    ("8×", "FASTER THAN LIMIT"),
    ("<200 MB", "RAM USAGE"),
    ("Zero", "EXTERNAL API CALLS"),
    ("✓ Passed", "VALIDATION SCRIPT"),
]

for i, (value, label) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 3.95)
    y = Inches(1.5 + row * 1.15)
    card = add_shape(slide8, x, y, Inches(3.6), Inches(1.0), LIGHT_BG)
    add_textbox(slide8, x, y + Inches(0.1), Inches(3.6), Inches(0.45),
                value, font_size=22, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, x, y + Inches(0.55), Inches(3.6), Inches(0.3),
                label, font_size=8, bold=True, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# Comparison table
table_y = Inches(3.95)
table_data = [
    ["Capability", "Our System", "Traditional ATS"],
    ["Matching Method", "Semantic cluster + TF-IDF + endorsements", "Exact keyword frequency"],
    ["Career Analysis", "Needleman-Wunsch alignment + PageRank pedigree", "Simple years-of-experience count"],
    ["Anti-Gaming", "Shannon entropy + multi-layer honeypot", "None"],
    ["Explainability", "Deterministic XAI rationale per candidate", "None / generic template"],
    ["Continuous Learning", "Recruiter feedback → perceptron weight update", "Manual rule tuning"],
]

rows, cols = len(table_data), 3
tbl = slide8.shapes.add_table(rows, cols, Inches(0.8), table_y, Inches(11.5), Inches(2.8)).table

tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(5.0)
tbl.columns[2].width = Inches(4.0)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.name = 'Calibri'
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF7, 0xFF)
            elif c == 1:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            else:
                paragraph.font.color.rgb = GRAY_700


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNOLOGIES USED
# ══════════════════════════════════════════════════════════════════
slide9 = make_content_slide("Technologies Used")

# Left side — Tech categories
tech_categories = [
    ("BACKEND / ENGINE", ["Python 3.8+", "multiprocessing", "json", "csv", "math", "re", "gzip", "argparse"], PURPLE),
    ("ALGORITHMS & TECHNIQUES", ["TF-IDF Weighting", "Needleman-Wunsch DP", "PageRank", "Shannon Entropy", "Exponential Decay", "Perceptron Learning"], BLUE),
    ("FRONTEND / DASHBOARD", ["HTML5", "Vanilla JavaScript", "CSS3", "LocalStorage API", "Zero Dependencies"], GREEN),
]

y_pos = 1.5
for cat_title, tags, tag_color in tech_categories:
    add_textbox(slide9, Inches(0.8), Inches(y_pos), Inches(5.5), Inches(0.25),
                cat_title, font_size=9, bold=True, color=GRAY_500)
    y_pos += 0.3
    # Tags as text (since python-pptx can't do styled chips easily)
    tag_text = "  •  ".join(tags)
    tag_box = add_textbox(slide9, Inches(0.8), Inches(y_pos), Inches(5.8), Inches(0.45),
                          tag_text, font_size=10, bold=False, color=tag_color)
    y_pos += 0.6

# Right side — Key design decisions
add_textbox(slide9, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.25),
            "KEY DESIGN DECISIONS", font_size=9, bold=True, color=GRAY_500)

add_bullet_textbox(slide9, Inches(6.8), Inches(1.85), Inches(5.5), Inches(4.5), [
    ("Pure Python stdlib:", "Zero external dependencies for the core engine — no pip install required. Guarantees sandbox compatibility."),
    ("Multiprocessing pool:", "Parallelizes scoring across all CPU cores with optimized chunk sizes for minimal IPC overhead."),
    ("Single-file frontend:", "The entire dashboard is one HTML file — no Node.js, no build tools, no server required."),
    ("JSON-driven config:", "All role requirements in one file — swap JD without code changes."),
], font_size=11)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — SUBMISSION ASSETS
# ══════════════════════════════════════════════════════════════════
slide10 = make_content_slide("Submission Assets")

assets = [
    ("📄", "team_submission.csv", "Final ranked output with 100 candidates: candidate_id, rank, score, reasoning. Passed the official validate_submission.py with zero errors."),
    ("⚙️", "run.py", "Main orchestrator script. Usage: python run.py --candidates candidates.jsonl --out team_submission.csv"),
    ("🧠", "src/", "Core engine modules: scorer.py (4-axis ranking), reasoning_engine.py (XAI), config.py (weights & JD rules loader)"),
    ("🌐", "frontend/index.html", "Interactive recruiter dashboard. Zero-dependency single-file web app with weight sliders, candidate cards, search, and active learning."),
    ("📋", "data/jd_rules.json", "Job description configuration: role, semantic clusters, experience bands, locations, company tiers. Fully customizable without code changes."),
    ("📝", "system_design.md", "Complete system architecture & design specification with Mermaid diagrams and algorithmic formulas."),
    ("🔗", "GitHub Repo", "github.com/Praveen-ing/redrob-ai-ranker  |  AI Tools: Cursor/Gemini AI for code architecture and pair programming"),
]

y_pos = 1.5
for (icon, name, desc) in assets:
    # Asset row background
    row_bg = add_shape(slide10, Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.65), LIGHT_BG)

    # Icon
    add_textbox(slide10, Inches(0.9), Inches(y_pos + 0.08), Inches(0.4), Inches(0.45),
                icon, font_size=14, alignment=PP_ALIGN.CENTER)

    # Name + description
    txBox = slide10.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.08), Inches(10.7), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_b = p.add_run()
    run_b.text = name + " — "
    run_b.font.size = Pt(11)
    run_b.font.bold = True
    run_b.font.color.rgb = BLACK
    run_b.font.name = 'Calibri'
    run_n = p.add_run()
    run_n.text = desc
    run_n.font.size = Pt(10)
    run_n.font.color.rgb = GRAY_700
    run_n.font.name = 'Calibri'

    y_pos += 0.75


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ══════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank_layout)

# Full dark background
add_shape(slide11, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x14, 0x0A, 0x2E))
add_shape(slide11, Inches(0), Inches(0), Inches(5), Inches(7.5), RGBColor(0x1A, 0x05, 0x33))
add_shape(slide11, Inches(8), Inches(2), Inches(5.333), Inches(5.5), RGBColor(0x4A, 0x19, 0x42))

# Logos
add_textbox(slide11, Inches(4), Inches(1.2), Inches(5.333), Inches(0.5),
            "● redrob  │  H2S", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Thank You
add_textbox(slide11, Inches(2), Inches(2.2), Inches(9.333), Inches(1.2),
            "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tagline
add_textbox(slide11, Inches(3), Inches(3.5), Inches(7.333), Inches(0.5),
            "Team TOPGUN — Building what next India runs on",
            font_size=14, bold=False, color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# Team members
add_textbox(slide11, Inches(2), Inches(4.5), Inches(9.333), Inches(0.4),
            "Nethavath Praveen  ·  Sri Abhinaya  ·  Botsa Radesh  ·  Sai Sameer Killamsetti",
            font_size=12, bold=False, color=GRAY_400, alignment=PP_ALIGN.CENTER)

# GitHub link
add_textbox(slide11, Inches(3), Inches(5.1), Inches(7.333), Inches(0.4),
            "github.com/Praveen-ing/redrob-ai-ranker",
            font_size=11, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Accent bar
add_accent_bar(slide11)


# ──────────────── SAVE ────────────────
output_path = "presentation_TOPGUN.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"   {len(prs.slides)} slides generated successfully.")
